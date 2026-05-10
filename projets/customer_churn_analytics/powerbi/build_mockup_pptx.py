"""
build_mockup_pptx.py — Genere le mockup PowerPoint vierge des 5 pages du dashboard
Customer Churn Analytics. Chaque slide = chrome (logo, navbar, footer) + cards
blanches vides en placeholder. Aucune donnee, aucun titre, aucun KPI value.

Charte Light + Orange chaleureux IvoirCom :
  - fond page    : #F9F9F8
  - cards        : #FFFFFF (bordure 0.5px #E5DDD5, ombre subtile)
  - primaire     : #E94E1B (logo, navbar active)
  - sidebar fonce: #A8350E (footer)
  - texte        : #2C2C2A / #6B7280

Lance :
    pip install python-pptx
    python build_mockup_pptx.py

Sortie :
  - mockup_customer_churn.pptx (5 slides, 1280x720 px)

Etapes suivantes :
  1) Ouvrir le PPTX, retoucher si besoin (positions, espacements)
  2) Enregistrer sous PNG : Fichier -> Enregistrer sous -> PNG (toutes les diapos)
  3) Renommer : bg-01-vue-executive.png ... bg-05-plan-action.png
  4) Importer comme arriere-plan dans Power BI Desktop (Format de la page -> Image)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt


# ============================================================
# Charte Light + Orange chaleureux
# ============================================================
COLOR_BG_PAGE     = RGBColor(0xF9, 0xF9, 0xF8)
COLOR_BG_CARD     = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BORDER      = RGBColor(0xE5, 0xDD, 0xD5)
COLOR_BORDER_DARK = RGBColor(0xD3, 0xC8, 0xBA)
COLOR_PRIMARY     = RGBColor(0xE9, 0x4E, 0x1B)
COLOR_PRIMARY_DK  = RGBColor(0xA8, 0x35, 0x0E)
COLOR_TEXT        = RGBColor(0x2C, 0x2C, 0x2A)
COLOR_TEXT_2      = RGBColor(0x6B, 0x72, 0x80)
COLOR_TEXT_3      = RGBColor(0x88, 0x87, 0x80)
COLOR_NAV_INACTIVE = RGBColor(0x88, 0x87, 0x80)

# ============================================================
# Layout 1280 x 720 pixels (a 96 DPI)
# 1 px = 9525 EMU
# ============================================================
PX = 9525  # EMU par pixel
SLIDE_W_PX = 1280
SLIDE_H_PX = 720

PADDING = 24

# Top bar
TOPBAR_Y = 16
TOPBAR_H = 50

# Title area
TITLE_Y = TOPBAR_Y + TOPBAR_H + 12
TITLE_H = 38

# Content area
CONTENT_Y = TITLE_Y + TITLE_H + 12
CONTENT_BOTTOM = SLIDE_H_PX - 36

# Footer
FOOTER_Y = SLIDE_H_PX - 28
FOOTER_H = 18

NAV_LABELS = [
    ("Vue Executive", 1),
    ("Segments", 2),
    ("Cohortes", 3),
    ("Reclamations", 4),
    ("Plan d'action", 5),
]


def emu(px: float) -> int:
    """Conversion pixels -> EMU."""
    return int(px * PX)


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.5,
             shadow=False) -> None:
    """Helper : ajoute un rectangle plein avec bordure optionnelle."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h)
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    if not shadow:
        # python-pptx n'expose pas l'ombre simplement -> placeholder
        try:
            shape.shadow.inherit = False
        except Exception:
            pass
    shape.text_frame.text = ""
    return shape


def add_rounded_card(slide, x, y, w, h) -> None:
    """Card blanche avec coins arrondis legers + bordure subtile."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h)
    )
    shape.adjustments[0] = 0.04  # rayon des coins ~5px
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG_CARD
    shape.line.color.rgb = COLOR_BORDER
    shape.line.width = Pt(0.5)
    shape.text_frame.text = ""
    return shape


def add_text(slide, x, y, w, h, text, size=11, bold=False,
             color=COLOR_TEXT, align="left", italic=False,
             font_name="Segoe UI") -> None:
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_logo_square(slide, x, y, size=36) -> None:
    """Carre orange avec lettre 'I' blanche."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(size), emu(size)
    )
    shape.adjustments[0] = 0.18
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "I"
    run.font.name = "Segoe UI"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def add_chrome(slide, active_idx: int) -> None:
    """Ajoute le chrome commun (fond, top bar, logo, navbar, footer) sur la slide.
    `active_idx` : 1..5, position de l'item navbar a marquer comme actif.
    """
    # Fond page
    add_rect(slide, 0, 0, SLIDE_W_PX, SLIDE_H_PX, fill=COLOR_BG_PAGE, line=None)

    # ---- TOP BAR ----
    top_x = PADDING
    add_logo_square(slide, top_x, TOPBAR_Y + 7, size=36)

    # Brand text
    add_text(slide, top_x + 48, TOPBAR_Y + 6, 220, 18,
             "IvoirCom", size=14, bold=True, color=COLOR_TEXT)
    add_text(slide, top_x + 48, TOPBAR_Y + 28, 240, 14,
             "Customer Churn Analytics", size=10, color=COLOR_TEXT_2)

    # Navbar centree
    nav_total_w = 580
    nav_x_start = (SLIDE_W_PX - nav_total_w) // 2
    nav_item_w = nav_total_w // 5
    for i, (label, idx) in enumerate(NAV_LABELS):
        is_active = idx == active_idx
        x = nav_x_start + i * nav_item_w
        color = COLOR_PRIMARY if is_active else COLOR_NAV_INACTIVE
        add_text(slide, x, TOPBAR_Y + 18, nav_item_w, 18,
                 label, size=11, bold=is_active, color=color, align="center")
        if is_active:
            # Soulignement orange 2px
            underline_w = 90
            underline_x = x + (nav_item_w - underline_w) // 2
            add_rect(slide, underline_x, TOPBAR_Y + 36, underline_w, 2,
                     fill=COLOR_PRIMARY, line=None)

    # Slicer placeholder (top-right)
    slicer_w = 130
    slicer_x = SLIDE_W_PX - PADDING - slicer_w
    slicer_y = TOPBAR_Y + 14
    slicer_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(slicer_x), emu(slicer_y), emu(slicer_w), emu(28)
    )
    slicer_shape.adjustments[0] = 0.3
    slicer_shape.fill.solid()
    slicer_shape.fill.fore_color.rgb = COLOR_BG_CARD
    slicer_shape.line.color.rgb = COLOR_BORDER
    slicer_shape.line.width = Pt(0.5)
    slicer_shape.text_frame.text = ""

    # Ligne de separation top bar
    add_rect(slide, PADDING, TOPBAR_Y + TOPBAR_H + 2,
             SLIDE_W_PX - 2 * PADDING, 0.5,
             fill=COLOR_BORDER, line=None)

    # ---- FOOTER ----
    add_rect(slide, PADDING, FOOTER_Y - 3,
             SLIDE_W_PX - 2 * PADDING, 0.5,
             fill=COLOR_BORDER, line=None)
    add_text(slide, PADDING, FOOTER_Y, 400, FOOTER_H,
             "DataProjectLab  -  IvoirCom  -  2024-2025",
             size=9, color=COLOR_TEXT_3)
    add_text(slide, SLIDE_W_PX - PADDING - 80, FOOTER_Y, 80, FOOTER_H,
             f"Page {active_idx} / 5", size=9, color=COLOR_TEXT_3,
             align="right")


# ============================================================
# Cartes placeholder par page
# ============================================================

def page_1_executive(slide) -> None:
    """Page 1 : Bandeau alerte + 4 KPI cards + 2 charts row + 1 chart row."""
    # Bandeau alerte (placeholder rouge pale)
    bandeau_y = CONTENT_Y
    bandeau_h = 36
    bandeau_w = SLIDE_W_PX - 2 * PADDING
    bandeau = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(PADDING), emu(bandeau_y), emu(bandeau_w), emu(bandeau_h)
    )
    bandeau.adjustments[0] = 0.1
    bandeau.fill.solid()
    bandeau.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xEA)  # rouge tres pale
    bandeau.line.color.rgb = RGBColor(0xE2, 0x4B, 0x4A)
    bandeau.line.width = Pt(0.5)
    bandeau.text_frame.text = ""

    # 4 KPI cards
    kpi_y = bandeau_y + bandeau_h + 12
    kpi_h = 100
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 3 * gap
    kpi_w = avail_w // 4
    for i in range(4):
        x = PADDING + i * (kpi_w + gap)
        add_rounded_card(slide, x, kpi_y, kpi_w, kpi_h)

    # 2 charts row : donut + line chart
    charts_y = kpi_y + kpi_h + 12
    charts_h = 200
    donut_w = (SLIDE_W_PX - 2 * PADDING - gap) // 3
    line_w = SLIDE_W_PX - 2 * PADDING - donut_w - gap
    add_rounded_card(slide, PADDING, charts_y, donut_w, charts_h)
    add_rounded_card(slide, PADDING + donut_w + gap, charts_y, line_w, charts_h)

    # Chart bottom : barres horizontales par offre
    bottom_y = charts_y + charts_h + 12
    bottom_h = CONTENT_BOTTOM - bottom_y
    add_rounded_card(slide, PADDING, bottom_y,
                     SLIDE_W_PX - 2 * PADDING, bottom_h)


def page_2_segments(slide) -> None:
    """Page 2 : 3 KPI + heatmap matrix + 2 small charts row."""
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    for i in range(3):
        x = PADDING + i * (kpi_w + gap)
        add_rounded_card(slide, x, kpi_y, kpi_w, kpi_h)

    # Heatmap matrix (large)
    heat_y = kpi_y + kpi_h + 12
    heat_h = 230
    add_rounded_card(slide, PADDING, heat_y,
                     SLIDE_W_PX - 2 * PADDING, heat_h)

    # 2 small charts (ville + age)
    small_y = heat_y + heat_h + 12
    small_h = CONTENT_BOTTOM - small_y
    small_w = (SLIDE_W_PX - 2 * PADDING - gap) // 2
    add_rounded_card(slide, PADDING, small_y, small_w, small_h)
    add_rounded_card(slide, PADDING + small_w + gap, small_y, small_w, small_h)


def page_3_cohortes(slide) -> None:
    """Page 3 : 3 KPI + heatmap retention 12x12 + table top cohortes."""
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    for i in range(3):
        x = PADDING + i * (kpi_w + gap)
        add_rounded_card(slide, x, kpi_y, kpi_w, kpi_h)

    # Heatmap retention (very wide)
    heat_y = kpi_y + kpi_h + 12
    heat_h = 280
    add_rounded_card(slide, PADDING, heat_y,
                     SLIDE_W_PX - 2 * PADDING, heat_h)

    # Table cohortes
    table_y = heat_y + heat_h + 12
    table_h = CONTENT_BOTTOM - table_y
    add_rounded_card(slide, PADDING, table_y,
                     SLIDE_W_PX - 2 * PADDING, table_h)


def page_4_reclamations(slide) -> None:
    """Page 4 : 3 KPI + 2 charts row + bandeau hero + table."""
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    for i in range(3):
        x = PADDING + i * (kpi_w + gap)
        add_rounded_card(slide, x, kpi_y, kpi_w, kpi_h)

    # 2 charts (type x statut + tickets actifs/churners)
    charts_y = kpi_y + kpi_h + 12
    charts_h = 170
    chart1_w = int((SLIDE_W_PX - 2 * PADDING - gap) * 0.55)
    chart2_w = SLIDE_W_PX - 2 * PADDING - chart1_w - gap
    add_rounded_card(slide, PADDING, charts_y, chart1_w, charts_h)
    add_rounded_card(slide, PADDING + chart1_w + gap, charts_y, chart2_w, charts_h)

    # Bandeau hero (signal phare)
    hero_y = charts_y + charts_h + 12
    hero_h = 80
    hero = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(PADDING), emu(hero_y),
        emu(SLIDE_W_PX - 2 * PADDING), emu(hero_h)
    )
    hero.adjustments[0] = 0.06
    hero.fill.solid()
    hero.fill.fore_color.rgb = RGBColor(0xFD, 0xEC, 0xEA)
    hero.line.color.rgb = RGBColor(0xE2, 0x4B, 0x4A)
    hero.line.width = Pt(1.5)
    hero.text_frame.text = ""

    # Table top clients
    table_y = hero_y + hero_h + 10
    table_h = CONTENT_BOTTOM - table_y
    add_rounded_card(slide, PADDING, table_y,
                     SLIDE_W_PX - 2 * PADDING, table_h)


def page_5_plan_action(slide) -> None:
    """Page 5 : 3 KPI + (donut + barres) + bandeau ARPU + 3 leviers."""
    # 3 KPI cards (border-left colore)
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    border_colors = [
        RGBColor(0x1D, 0x9E, 0x75),  # Champions vert
        RGBColor(0xEF, 0x9F, 0x27),  # At Risk orange
        RGBColor(0xE2, 0x4B, 0x4A),  # Lost rouge
    ]
    for i in range(3):
        x = PADDING + i * (kpi_w + gap)
        add_rounded_card(slide, x, kpi_y, kpi_w, kpi_h)
        # bordure gauche coloree (rectangle plein 3px)
        add_rect(slide, x, kpi_y, 3, kpi_h, fill=border_colors[i], line=None)

    # Donut + barres
    charts_y = kpi_y + kpi_h + 12
    charts_h = 180
    donut_w = (SLIDE_W_PX - 2 * PADDING - gap) // 2 - 50
    bar_w = SLIDE_W_PX - 2 * PADDING - donut_w - gap
    add_rounded_card(slide, PADDING, charts_y, donut_w, charts_h)
    add_rounded_card(slide, PADDING + donut_w + gap, charts_y, bar_w, charts_h)

    # Bandeau ARPU expose (orange pale)
    arpu_y = charts_y + charts_h + 10
    arpu_h = 60
    arpu = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(PADDING), emu(arpu_y),
        emu(SLIDE_W_PX - 2 * PADDING), emu(arpu_h)
    )
    arpu.adjustments[0] = 0.08
    arpu.fill.solid()
    arpu.fill.fore_color.rgb = RGBColor(0xFF, 0xF4, 0xED)
    arpu.line.color.rgb = COLOR_PRIMARY
    arpu.line.width = Pt(1.0)
    arpu.text_frame.text = ""

    # 3 leviers card (bas)
    levers_y = arpu_y + arpu_h + 10
    levers_h = CONTENT_BOTTOM - levers_y
    add_rounded_card(slide, PADDING, levers_y,
                     SLIDE_W_PX - 2 * PADDING, levers_h)


# ============================================================
# Main
# ============================================================

PAGE_BUILDERS = [
    ("Vue Executive", page_1_executive),
    ("Segments", page_2_segments),
    ("Cohortes", page_3_cohortes),
    ("Reclamations", page_4_reclamations),
    ("Plan d'action", page_5_plan_action),
]


def main() -> None:
    out_path = Path(__file__).parent / "mockup_customer_churn.pptx"

    prs = Presentation()
    # Slide size 1280 x 720 px
    prs.slide_width = emu(SLIDE_W_PX)
    prs.slide_height = emu(SLIDE_H_PX)

    blank_layout = prs.slide_layouts[6]  # layout vide

    for idx, (name, builder) in enumerate(PAGE_BUILDERS, start=1):
        slide = prs.slides.add_slide(blank_layout)
        add_chrome(slide, active_idx=idx)
        builder(slide)

    prs.save(out_path)
    size_kb = out_path.stat().st_size / 1024
    print(f"[OK] Mockup PPTX genere : {out_path}")
    print(f"     5 slides 1280x720 px  -  {size_kb:.1f} Ko")
    print()
    print("Etapes suivantes :")
    print("  1) Ouvrir le PPTX, retoucher si besoin")
    print("  2) Fichier -> Enregistrer sous -> PNG -> Toutes les diapositives")
    print("     (Options PowerPoint : ExportBitmapResolution = 150 via regedit)")
    print("  3) Renommer les PNG :")
    print("     bg-01-vue-executive.png")
    print("     bg-02-segments.png")
    print("     bg-03-cohortes.png")
    print("     bg-04-reclamations.png")
    print("     bg-05-plan-action.png")
    print("  4) Importer comme arriere-plan dans Power BI Desktop")
    print("     (Format de la page -> Image -> Adapter -> Transparence 0%)")


if __name__ == "__main__":
    main()
