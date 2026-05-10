"""
build_mockup_pptx_dark.py — Mockup PowerPoint vierge, version Dark + icones.

Charte Dark + Orange chaleureux IvoirCom :
  - fond page    : #0B1421
  - cards        : #1A2436 (bordure 0.5px #2A3548)
  - primaire     : #E94E1B (logo, navbar active, accents)
  - texte        : #FFFFFF / #B0B8C7 / #6B7B95

Iconographie :
  - Logo IvoirCom : carre orange + glyphe stylise smartphone
  - Navbar : 5 icones SVG composees a partir de formes primitives
    (cible, barres, reseau, ticket, eclair)
  - Slicer : icone calendrier
  - KPI cards : icone fantome en haut-gauche (placeholder)

Lance :
    pip install python-pptx
    python build_mockup_pptx_dark.py

Sortie : mockup_customer_churn_dark.pptx (5 slides 1280x720 px)
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


# ============================================================
# Charte Dark + Orange chaleureux
# ============================================================
COLOR_BG_PAGE     = RGBColor(0x0B, 0x14, 0x21)
COLOR_BG_CARD     = RGBColor(0x1A, 0x24, 0x36)
COLOR_BG_CARD_HI  = RGBColor(0x23, 0x2F, 0x45)
COLOR_BORDER      = RGBColor(0x2A, 0x35, 0x48)
COLOR_PRIMARY     = RGBColor(0xE9, 0x4E, 0x1B)
COLOR_PRIMARY_DK  = RGBColor(0xA8, 0x35, 0x0E)
COLOR_TEXT        = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT_2      = RGBColor(0xB0, 0xB8, 0xC7)
COLOR_TEXT_3      = RGBColor(0x6B, 0x7B, 0x95)
COLOR_NAV_INACTIVE = RGBColor(0x6B, 0x7B, 0x95)
COLOR_GREEN       = RGBColor(0x10, 0xD9, 0xA3)
COLOR_RED         = RGBColor(0xFF, 0x4D, 0x6D)
COLOR_WARN        = RGBColor(0xFF, 0xB5, 0x47)

# Bandeau alerte (page 1) — fond rouge tres fonce + bordure rouge corail
COLOR_ALERT_BG    = RGBColor(0x3D, 0x14, 0x21)

# ============================================================
# Layout 1280 x 720 pixels (96 DPI)
# ============================================================
PX = 9525
SLIDE_W_PX = 1280
SLIDE_H_PX = 720
PADDING = 24
TOPBAR_Y = 16
TOPBAR_H = 50
TITLE_Y = TOPBAR_Y + TOPBAR_H + 12
TITLE_H = 38
CONTENT_Y = TITLE_Y + TITLE_H + 12
CONTENT_BOTTOM = SLIDE_H_PX - 36
FOOTER_Y = SLIDE_H_PX - 28
FOOTER_H = 18

NAV_LABELS = [
    ("Vue Executive", 1, "target"),
    ("Segments",      2, "bar"),
    ("Cohortes",      3, "network"),
    ("Reclamations",  4, "ticket"),
    ("Plan d'action", 5, "bolt"),
]


def emu(px: float) -> int:
    return int(px * PX)


# ============================================================
# Helpers de base
# ============================================================

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.5):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h)
    )
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w)
    shape.text_frame.text = ""
    return shape


def add_rounded_card(slide, x, y, w, h, fill=COLOR_BG_CARD,
                     border=COLOR_BORDER, radius=0.04) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h)
    )
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.5)
    shape.text_frame.text = ""
    return shape


def add_text(slide, x, y, w, h, text, size=11, bold=False,
             color=COLOR_TEXT, align="left", italic=False,
             font_name="Segoe UI"):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                   "right": PP_ALIGN.RIGHT}[align]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def add_oval(slide, x, y, size, fill=None, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(x), emu(y), emu(size), emu(size))
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.text_frame.text = ""
    return s


# ============================================================
# Bibliotheque d'icones (composees de primitives)
# ============================================================

def icon_target(slide, x, y, size=14, color=COLOR_TEXT_2):
    """Cible : 2 cercles concentriques + point central."""
    add_oval(slide, x, y, size, line=color, line_w=1.2)
    inner = size * 0.45
    ix = x + (size - inner) / 2
    iy = y + (size - inner) / 2
    add_oval(slide, ix, iy, inner, line=color, line_w=1.0)
    dot = size * 0.18
    dx = x + (size - dot) / 2
    dy = y + (size - dot) / 2
    add_oval(slide, dx, dy, dot, fill=color)


def icon_bar(slide, x, y, size=14, color=COLOR_TEXT_2):
    """3 barres verticales croissantes."""
    w = size * 0.22
    gap = size * 0.1
    heights = [size * 0.45, size * 0.7, size * 0.95]
    for i, h in enumerate(heights):
        bx = x + i * (w + gap) + size * 0.05
        by = y + size - h
        add_rect(slide, bx, by, w, h, fill=color)


def icon_network(slide, x, y, size=14, color=COLOR_TEXT_2):
    """3 noeuds en triangle."""
    r = size * 0.22
    positions = [(0, 0), (size - r * 2, 0), (size / 2 - r, size - r * 2)]
    for px, py in positions:
        add_oval(slide, x + px, y + py, r * 2, fill=color)


def icon_ticket(slide, x, y, size=14, color=COLOR_TEXT_2):
    """Rectangle simulant un ticket (avec entaille via trait diagonal)."""
    h = size * 0.7
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(x), emu(y + (size - h) / 2), emu(size), emu(h)
    )
    s.adjustments[0] = 0.18
    s.fill.background()
    s.line.color.rgb = color
    s.line.width = Pt(1.2)
    s.text_frame.text = ""
    # ligne pointillee verticale au centre
    line_x = x + size / 2
    add_rect(slide, line_x, y + (size - h) / 2 + 2, 0.5, h - 4, fill=color)


def icon_bolt(slide, x, y, size=14, color=COLOR_TEXT_2):
    """Eclair (MSO_SHAPE built-in)."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.LIGHTNING_BOLT,
        emu(x + size * 0.15), emu(y), emu(size * 0.7), emu(size)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.text_frame.text = ""


def icon_calendar(slide, x, y, size=12, color=COLOR_TEXT_2):
    """Calendrier : rectangle + barre du haut + 2 anneaux."""
    add_rect(slide, x, y + size * 0.15, size, size * 0.85,
             fill=None, line=color, line_w=1.0)
    # barre du haut (header)
    add_rect(slide, x, y + size * 0.15, size, size * 0.2, fill=color)
    # 2 petits anneaux
    ring = size * 0.1
    add_rect(slide, x + size * 0.2, y, ring, size * 0.2, fill=color)
    add_rect(slide, x + size * 0.7, y, ring, size * 0.2, fill=color)


def icon_filter(slide, x, y, size=12, color=COLOR_TEXT_2):
    """Entonnoir simple : trapeze + rectangle."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        emu(x), emu(y), emu(size), emu(size * 0.6)
    )
    s.rotation = 180  # pointe vers le bas
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    # tige
    add_rect(slide, x + size * 0.4, y + size * 0.55, size * 0.2, size * 0.45, fill=color)


def icon_pin(slide, x, y, size=14, color=COLOR_TEXT_2):
    """Pin (goutte stylisee)."""
    add_oval(slide, x + size * 0.1, y, size * 0.8, fill=color)
    # petite pointe via triangle
    s = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        emu(x + size * 0.3), emu(y + size * 0.55), emu(size * 0.4), emu(size * 0.45)
    )
    s.rotation = 180
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


def icon_phone(slide, x, y, size=14, color=COLOR_TEXT):
    """Smartphone : rectangle arrondi + petit trait pour le bouton."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(x + size * 0.2), emu(y), emu(size * 0.6), emu(size)
    )
    s.adjustments[0] = 0.18
    s.fill.background()
    s.line.color.rgb = color
    s.line.width = Pt(1.3)
    s.text_frame.text = ""
    # trait bouton home
    btn_w = size * 0.2
    btn_h = size * 0.05
    add_rect(slide, x + (size - btn_w) / 2, y + size * 0.85, btn_w, btn_h, fill=color)


def icon_chevron(slide, x, y, size=10, color=COLOR_TEXT_2):
    """Chevron pointant vers le bas (MSO triangle)."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        emu(x), emu(y), emu(size), emu(size * 0.6)
    )
    s.rotation = 180
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()


# Dispatch icone par nom
ICONS = {
    "target":   icon_target,
    "bar":      icon_bar,
    "network":  icon_network,
    "ticket":   icon_ticket,
    "bolt":     icon_bolt,
    "calendar": icon_calendar,
    "filter":   icon_filter,
    "pin":      icon_pin,
    "phone":    icon_phone,
    "chevron":  icon_chevron,
}


def add_icon(slide, name, x, y, size=14, color=COLOR_TEXT_2):
    """Dispatcher d'icone."""
    fn = ICONS.get(name)
    if fn:
        fn(slide, x, y, size, color)


# ============================================================
# Chrome (top bar + footer) commun
# ============================================================

def add_logo(slide, x, y, size=36):
    """Carre orange avec icone phone blanche."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(size), emu(size)
    )
    s.adjustments[0] = 0.18
    s.fill.solid()
    s.fill.fore_color.rgb = COLOR_PRIMARY
    s.line.fill.background()
    s.text_frame.text = ""
    # Icone phone centre dans le carre
    icon_size = size * 0.55
    ix = x + (size - icon_size) / 2
    iy = y + (size - icon_size) / 2
    icon_phone(slide, ix, iy, size=icon_size, color=COLOR_TEXT)


def add_chrome(slide, active_idx: int) -> None:
    """Top bar (logo + brand + navbar + slicer) + footer."""
    # Fond page
    add_rect(slide, 0, 0, SLIDE_W_PX, SLIDE_H_PX, fill=COLOR_BG_PAGE)

    # Logo + Brand
    add_logo(slide, PADDING, TOPBAR_Y + 7, size=36)
    add_text(slide, PADDING + 48, TOPBAR_Y + 6, 240, 18,
             "IvoirCom", size=14, bold=True, color=COLOR_TEXT)
    add_text(slide, PADDING + 48, TOPBAR_Y + 28, 260, 14,
             "Customer Churn Analytics", size=10, color=COLOR_TEXT_2)

    # Navbar centree avec icones
    nav_total_w = 620
    nav_x_start = (SLIDE_W_PX - nav_total_w) // 2
    nav_item_w = nav_total_w // 5
    for i, (label, idx, icon_name) in enumerate(NAV_LABELS):
        is_active = idx == active_idx
        x_item = nav_x_start + i * nav_item_w
        color = COLOR_PRIMARY if is_active else COLOR_NAV_INACTIVE

        # icone (10px) a gauche du texte
        icon_size = 11
        icon_x = x_item + 8
        icon_y = TOPBAR_Y + 22
        add_icon(slide, icon_name, icon_x, icon_y, size=icon_size, color=color)

        # texte
        text_x = icon_x + icon_size + 5
        text_w = nav_item_w - (icon_size + 13)
        add_text(slide, text_x, TOPBAR_Y + 19, text_w, 16,
                 label, size=11, bold=is_active, color=color, align="left")

        if is_active:
            underline_x = x_item + 8
            underline_w = nav_item_w - 16
            add_rect(slide, underline_x, TOPBAR_Y + 39, underline_w, 2,
                     fill=COLOR_PRIMARY)

    # Slicer en haut-droite (avec icone calendar/filter selon page)
    slicer_w = 140
    slicer_h = 28
    slicer_x = SLIDE_W_PX - PADDING - slicer_w
    slicer_y = TOPBAR_Y + 14
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(slicer_x), emu(slicer_y), emu(slicer_w), emu(slicer_h)
    )
    s.adjustments[0] = 0.3
    s.fill.solid()
    s.fill.fore_color.rgb = COLOR_BG_CARD
    s.line.color.rgb = COLOR_BORDER
    s.line.width = Pt(0.5)
    s.text_frame.text = ""
    # icone slicer (calendar pour pages 1/3, pin pour page 2, filter pour pages 4/5)
    slicer_icon = "calendar" if active_idx in (1, 3) else "pin" if active_idx == 2 else "filter"
    add_icon(slide, slicer_icon, slicer_x + 10, slicer_y + 8, size=12, color=COLOR_TEXT_2)
    # chevron a droite
    add_icon(slide, "chevron", slicer_x + slicer_w - 16, slicer_y + 11, size=8, color=COLOR_TEXT_2)

    # Ligne separation top bar
    add_rect(slide, PADDING, TOPBAR_Y + TOPBAR_H + 2,
             SLIDE_W_PX - 2 * PADDING, 0.5, fill=COLOR_BORDER)

    # Footer
    add_rect(slide, PADDING, FOOTER_Y - 3,
             SLIDE_W_PX - 2 * PADDING, 0.5, fill=COLOR_BORDER)
    add_text(slide, PADDING, FOOTER_Y, 400, FOOTER_H,
             "DataProjectLab  -  IvoirCom  -  2024-2025",
             size=9, color=COLOR_TEXT_3)
    add_text(slide, SLIDE_W_PX - PADDING - 80, FOOTER_Y, 80, FOOTER_H,
             f"Page {active_idx} / 5", size=9, color=COLOR_TEXT_3, align="right")


def add_kpi_card_with_icon(slide, x, y, w, h, icon_name, icon_color):
    """KPI card placeholder avec icone fantome en haut-gauche."""
    add_rounded_card(slide, x, y, w, h)
    # Pastille icone (background card-hi + icone)
    icn_size = 22
    icn_x = x + 12
    icn_y = y + 12
    pastille = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(icn_x), emu(icn_y), emu(icn_size), emu(icn_size)
    )
    pastille.adjustments[0] = 0.22
    pastille.fill.solid()
    pastille.fill.fore_color.rgb = COLOR_BG_CARD_HI
    pastille.line.fill.background()
    pastille.text_frame.text = ""
    # Icone elle-meme
    inner_size = 13
    inner_x = icn_x + (icn_size - inner_size) / 2
    inner_y = icn_y + (icn_size - inner_size) / 2
    add_icon(slide, icon_name, inner_x, inner_y, size=inner_size, color=icon_color)


# ============================================================
# Page builders
# ============================================================

def page_1_executive(slide):
    # Bandeau alerte
    bandeau_y = CONTENT_Y
    bandeau_h = 36
    bandeau_w = SLIDE_W_PX - 2 * PADDING
    b = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(PADDING), emu(bandeau_y), emu(bandeau_w), emu(bandeau_h)
    )
    b.adjustments[0] = 0.1
    b.fill.solid()
    b.fill.fore_color.rgb = COLOR_ALERT_BG
    b.line.color.rgb = COLOR_RED
    b.line.width = Pt(0.5)
    b.text_frame.text = ""
    # Bordure gauche 3px rouge
    add_rect(slide, PADDING, bandeau_y, 3, bandeau_h, fill=COLOR_RED)

    # 4 KPI cards
    kpi_y = bandeau_y + bandeau_h + 12
    kpi_h = 100
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 3 * gap
    kpi_w = avail_w // 4
    kpi_icons = [
        ("network", COLOR_PRIMARY),  # users-like (Total Clients)
        ("filter",  COLOR_RED),       # warning (Taux Churn)
        ("target",  COLOR_GREEN),     # money (ARPU)
        ("ticket",  COLOR_WARN),      # tickets (Reclamations)
    ]
    for i, (icon, color) in enumerate(kpi_icons):
        x = PADDING + i * (kpi_w + gap)
        add_kpi_card_with_icon(slide, x, kpi_y, kpi_w, kpi_h, icon, color)

    # Charts row
    charts_y = kpi_y + kpi_h + 12
    charts_h = 200
    donut_w = (SLIDE_W_PX - 2 * PADDING - gap) // 3
    line_w = SLIDE_W_PX - 2 * PADDING - donut_w - gap
    add_rounded_card(slide, PADDING, charts_y, donut_w, charts_h)
    add_rounded_card(slide, PADDING + donut_w + gap, charts_y, line_w, charts_h)

    # Chart bottom
    bottom_y = charts_y + charts_h + 12
    bottom_h = CONTENT_BOTTOM - bottom_y
    add_rounded_card(slide, PADDING, bottom_y, SLIDE_W_PX - 2 * PADDING, bottom_h)


def page_2_segments(slide):
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    kpi_icons = [
        ("filter", COLOR_RED),
        ("pin",    RGBColor(0x8B, 0x7D, 0xFF)),
        ("target", COLOR_PRIMARY),
    ]
    for i, (icon, color) in enumerate(kpi_icons):
        x = PADDING + i * (kpi_w + gap)
        add_kpi_card_with_icon(slide, x, kpi_y, kpi_w, kpi_h, icon, color)

    # Heatmap large
    heat_y = kpi_y + kpi_h + 12
    heat_h = 230
    add_rounded_card(slide, PADDING, heat_y, SLIDE_W_PX - 2 * PADDING, heat_h)

    # 2 small charts
    small_y = heat_y + heat_h + 12
    small_h = CONTENT_BOTTOM - small_y
    small_w = (SLIDE_W_PX - 2 * PADDING - gap) // 2
    add_rounded_card(slide, PADDING, small_y, small_w, small_h)
    add_rounded_card(slide, PADDING + small_w + gap, small_y, small_w, small_h)


def page_3_cohortes(slide):
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    kpi_icons = [
        ("target", COLOR_GREEN),
        ("target", COLOR_RED),
        ("network", COLOR_PRIMARY),
    ]
    for i, (icon, color) in enumerate(kpi_icons):
        x = PADDING + i * (kpi_w + gap)
        add_kpi_card_with_icon(slide, x, kpi_y, kpi_w, kpi_h, icon, color)

    heat_y = kpi_y + kpi_h + 12
    heat_h = 280
    add_rounded_card(slide, PADDING, heat_y, SLIDE_W_PX - 2 * PADDING, heat_h)

    table_y = heat_y + heat_h + 12
    table_h = CONTENT_BOTTOM - table_y
    add_rounded_card(slide, PADDING, table_y, SLIDE_W_PX - 2 * PADDING, table_h)


def page_4_reclamations(slide):
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    kpi_icons = [
        ("ticket", COLOR_WARN),
        ("filter", COLOR_RED),
        ("target", RGBColor(0x8B, 0x7D, 0xFF)),
    ]
    for i, (icon, color) in enumerate(kpi_icons):
        x = PADDING + i * (kpi_w + gap)
        add_kpi_card_with_icon(slide, x, kpi_y, kpi_w, kpi_h, icon, color)

    charts_y = kpi_y + kpi_h + 12
    charts_h = 170
    chart1_w = int((SLIDE_W_PX - 2 * PADDING - gap) * 0.55)
    chart2_w = SLIDE_W_PX - 2 * PADDING - chart1_w - gap
    add_rounded_card(slide, PADDING, charts_y, chart1_w, charts_h)
    add_rounded_card(slide, PADDING + chart1_w + gap, charts_y, chart2_w, charts_h)

    # Bandeau hero
    hero_y = charts_y + charts_h + 12
    hero_h = 80
    hero = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(PADDING), emu(hero_y),
        emu(SLIDE_W_PX - 2 * PADDING), emu(hero_h)
    )
    hero.adjustments[0] = 0.06
    hero.fill.solid()
    hero.fill.fore_color.rgb = COLOR_ALERT_BG
    hero.line.color.rgb = COLOR_RED
    hero.line.width = Pt(1.5)
    hero.text_frame.text = ""
    # Bord gauche 3px
    add_rect(slide, PADDING, hero_y, 3, hero_h, fill=COLOR_RED)

    table_y = hero_y + hero_h + 10
    table_h = CONTENT_BOTTOM - table_y
    add_rounded_card(slide, PADDING, table_y, SLIDE_W_PX - 2 * PADDING, table_h)


def page_5_plan_action(slide):
    kpi_y = CONTENT_Y
    kpi_h = 90
    gap = 12
    avail_w = SLIDE_W_PX - 2 * PADDING - 2 * gap
    kpi_w = avail_w // 3
    kpi_icons = [
        ("target",  COLOR_GREEN),
        ("filter",  COLOR_WARN),
        ("filter",  COLOR_RED),
    ]
    border_colors = [COLOR_GREEN, COLOR_WARN, COLOR_RED]
    for i, ((icon, color), border) in enumerate(zip(kpi_icons, border_colors)):
        x = PADDING + i * (kpi_w + gap)
        add_kpi_card_with_icon(slide, x, kpi_y, kpi_w, kpi_h, icon, color)
        # Bordure gauche 2px
        add_rect(slide, x, kpi_y, 2, kpi_h, fill=border)

    charts_y = kpi_y + kpi_h + 12
    charts_h = 180
    donut_w = (SLIDE_W_PX - 2 * PADDING - gap) // 2 - 50
    bar_w = SLIDE_W_PX - 2 * PADDING - donut_w - gap
    add_rounded_card(slide, PADDING, charts_y, donut_w, charts_h)
    add_rounded_card(slide, PADDING + donut_w + gap, charts_y, bar_w, charts_h)

    arpu_y = charts_y + charts_h + 10
    arpu_h = 60
    arpu = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        emu(PADDING), emu(arpu_y),
        emu(SLIDE_W_PX - 2 * PADDING), emu(arpu_h)
    )
    arpu.adjustments[0] = 0.08
    arpu.fill.solid()
    arpu.fill.fore_color.rgb = RGBColor(0x2A, 0x18, 0x10)
    arpu.line.color.rgb = COLOR_PRIMARY
    arpu.line.width = Pt(1.0)
    arpu.text_frame.text = ""

    levers_y = arpu_y + arpu_h + 10
    levers_h = CONTENT_BOTTOM - levers_y
    add_rounded_card(slide, PADDING, levers_y, SLIDE_W_PX - 2 * PADDING, levers_h)


# ============================================================
# Main
# ============================================================

PAGE_BUILDERS = [
    ("Vue Executive", page_1_executive),
    ("Segments",      page_2_segments),
    ("Cohortes",      page_3_cohortes),
    ("Reclamations",  page_4_reclamations),
    ("Plan d'action", page_5_plan_action),
]


def main() -> None:
    out_path = Path(__file__).parent / "mockup_customer_churn_dark.pptx"

    prs = Presentation()
    prs.slide_width = emu(SLIDE_W_PX)
    prs.slide_height = emu(SLIDE_H_PX)

    blank_layout = prs.slide_layouts[6]

    for idx, (name, builder) in enumerate(PAGE_BUILDERS, start=1):
        slide = prs.slides.add_slide(blank_layout)
        add_chrome(slide, active_idx=idx)
        builder(slide)

    prs.save(out_path)
    size_kb = out_path.stat().st_size / 1024
    print(f"[OK] Mockup PPTX Dark genere : {out_path}")
    print(f"     5 slides 1280x720 px  -  {size_kb:.1f} Ko")
    print()
    print("Iconographie :")
    print("  - Logo : carre orange + glyphe smartphone")
    print("  - Navbar : 5 icones (target, bar, network, ticket, bolt)")
    print("  - Slicer : icone calendar / pin / filter selon page")
    print("  - KPI cards : pastille icone fantome haut-gauche")
    print()
    print("Etapes suivantes :")
    print("  1) Ouvrir le PPTX, retoucher si besoin")
    print("  2) Activer ExportBitmapResolution = 150 via regedit")
    print("  3) Fichier -> Enregistrer sous -> PNG -> Toutes les diapositives")
    print("  4) Renommer : bg-01 a bg-05 (vue-executive, segments, cohortes, reclamations, plan-action)")
    print("  5) Importer comme arriere-plan dans Power BI Desktop")


if __name__ == "__main__":
    main()
