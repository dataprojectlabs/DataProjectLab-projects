"""
build_synthese_edutrack.py
Genere la presentation de synthese EduTrack en .pptx (10 slides).
Focus : insights business, pas recommandations generiques.

Sortie : synthese_edutrack.pptx (meme dossier)

Prerequis :
    pip install python-pptx requests

Usage :
    python build_synthese_edutrack.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pathlib import Path
import urllib.request

# ============================================================
# Configuration
# ============================================================

ROOT = Path(__file__).parent
OUT = ROOT / "synthese_edutrack.pptx"
LOGO_URL = "https://raw.githubusercontent.com/dataprojectlabs/DataProjectLab-projects/refs/heads/main/media/logo_dataprojectlab.png"
LOGO_CACHE = ROOT / ".cache_logo.png"

# Charte EduTrack — palette violet/vert sur navy
NAVY = RGBColor(0x1E, 0x3A, 0x5F)
NAVY_DARK = RGBColor(0x0F, 0x1F, 0x3A)
VIOLET = RGBColor(0x53, 0x4A, 0xB7)
VIOLET_LIGHT = RGBColor(0xEE, 0xED, 0xFE)
VERT = RGBColor(0x1F, 0xA6, 0x7D)
VERT_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
ROUGE = RGBColor(0xE5, 0x49, 0x4D)
ROUGE_LIGHT = RGBColor(0xFE, 0xE2, 0xE2)
ORANGE = RGBColor(0xF2, 0xA9, 0x3B)
ORANGE_LIGHT = RGBColor(0xFE, 0xF3, 0xC7)
GRIS = RGBColor(0x6B, 0x72, 0x80)
GRIS_LIGHT = RGBColor(0xE5, 0xE7, 0xEB)
BLANC = RGBColor(0xFF, 0xFF, 0xFF)
TEXTE = RGBColor(0x11, 0x18, 0x27)
FOND_PAGE = RGBColor(0xF9, 0xFA, 0xFB)


def download_logo():
    if not LOGO_CACHE.exists():
        urllib.request.urlretrieve(LOGO_URL, LOGO_CACHE)
    return LOGO_CACHE


# ============================================================
# Helpers
# ============================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK_LAYOUT = prs.slide_layouts[6]  # Blank


def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_rounded(slide, x, y, w, h, fill_rgb, line_rgb=None, radius=0.05):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, italic=False,
             color=TEXTE, font="Calibri", align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0)
    tf.margin_right = Pt(0)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tx


def add_multi(slide, x, y, w, h, lines, anchor=MSO_ANCHOR.TOP):
    """lines = [(text, size, bold, color, font, italic), ...]"""
    tx = slide.shapes.add_textbox(x, y, w, h)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    tf.margin_top = Pt(0); tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        if len(line) == 6:
            text, size, bold, color, font, italic = line
        else:
            text, size, bold, color, font = line
            italic = False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = text
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tx


def add_footer(slide, num, total=10):
    # Bandeau footer discret
    add_text(slide, Inches(0.5), Inches(7.15),
             Inches(8), Inches(0.25),
             "DataProjectLab — Academy   ·   EduTrack Analytics — Synthèse exécutive",
             size=8, color=GRIS, italic=True)
    add_text(slide, Inches(11.5), Inches(7.15),
             Inches(1.4), Inches(0.25),
             f"{num} / {total}",
             size=8, color=GRIS, align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 1 — COUVERTURE
# ============================================================

s = prs.slides.add_slide(BLANK_LAYOUT)
# Fond gradient simulé : 2 rectangles
add_rect(s, 0, 0, SW, SH, NAVY_DARK)
# Cercles décoratifs violets
c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(0.5), Inches(4.5), Inches(4.5))
c.fill.solid(); c.fill.fore_color.rgb = VIOLET
c.line.fill.background(); c.shadow.inherit = False
c.fill.transparency = 0  # pas dispo, on simule avec couleur foncée
# Petit cercle violet
c2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11), Inches(4.5), Inches(2.5), Inches(2.5))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x3C, 0x34, 0x89)
c2.line.fill.background(); c2.shadow.inherit = False

# Tag pill
pill = add_rounded(s, Inches(0.7), Inches(0.7), Inches(2.4), Inches(0.4),
                   VIOLET, radius=0.4)
add_text(s, Inches(0.7), Inches(0.72), Inches(2.4), Inches(0.4),
         "DATAPROJECTLAB · ACADEMY",
         size=9, bold=True, color=BLANC, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

# Titre principal
add_text(s, Inches(0.7), Inches(1.7), Inches(11), Inches(1.5),
         "EduTrack Analytics",
         size=64, bold=True, color=BLANC, font="Georgia")

add_text(s, Inches(0.7), Inches(2.9), Inches(11), Inches(1.2),
         "Synthèse exécutive",
         size=44, bold=False, color=VIOLET_LIGHT, font="Georgia", italic=True)

add_text(s, Inches(0.7), Inches(4.0), Inches(10), Inches(0.4),
         "Le pilotage data de DataProjectLab Academy",
         size=14, italic=True, color=RGBColor(0xCB, 0xD5, 0xE0))

# Trait fin
add_rect(s, Inches(0.7), Inches(4.5), Inches(0.8), Pt(2), VIOLET_LIGHT)

# 4 KPI bottom
kpi_y = Inches(5.6)
kpi_w = Inches(2.5)
kpi_h = Inches(1.0)
kpis_cover = [
    ("4 500", "Apprenants"),
    ("6 456", "Inscriptions"),
    ("33,3 %", "Complétion"),
    ("778M FCFA", "Revenu"),
]
for i, (val, lbl) in enumerate(kpis_cover):
    x = Inches(0.7 + i * 2.7)
    add_text(s, x, kpi_y, kpi_w, Inches(0.7), val,
             size=32, bold=True, color=BLANC, font="Georgia")
    add_text(s, x, Inches(6.4), kpi_w, Inches(0.4), lbl,
             size=12, color=VIOLET_LIGHT, italic=True)

# Date
add_text(s, Inches(11.2), Inches(6.9), Inches(1.8), Inches(0.4),
         "Mai 2026  ·  v1.0",
         size=10, color=RGBColor(0x9A, 0xB6, 0xE0), align=PP_ALIGN.RIGHT)


# ============================================================
# SLIDE 2 — L'ENJEU EN 1 PHRASE
# ============================================================

s = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s, 0, 0, SW, SH, FOND_PAGE)

# Tag
add_text(s, Inches(0.7), Inches(0.5), Inches(3), Inches(0.3),
         "01  ·  L'ENJEU",
         size=10, bold=True, color=VIOLET)

add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "EduTrack croît, mais perd 1 apprenant sur 6.",
         size=42, bold=True, color=NAVY, font="Georgia")

add_text(s, Inches(0.7), Inches(2.0), Inches(12), Inches(0.6),
         "Et le modèle ML signale un risque de décrochage massif sur les apprenants en cours.",
         size=18, italic=True, color=GRIS)

# 3 grands chiffres
y_big = Inches(3.4)
add_rounded(s, Inches(0.7), y_big, Inches(4.0), Inches(2.8), BLANC,
            line_rgb=GRIS_LIGHT, radius=0.05)
add_rect(s, Inches(0.7), y_big, Inches(4.0), Pt(4), VIOLET)
add_text(s, Inches(1.0), Inches(3.6), Inches(3.7), Inches(0.4),
         "INSCRIPTIONS",
         size=10, bold=True, color=VIOLET)
add_text(s, Inches(1.0), Inches(4.0), Inches(3.7), Inches(1.0),
         "+21 %",
         size=56, bold=True, color=NAVY, font="Georgia")
add_text(s, Inches(1.0), Inches(5.0), Inches(3.7), Inches(0.4),
         "vs N-1  ·  6 456 inscriptions",
         size=12, color=GRIS)
add_text(s, Inches(1.0), Inches(5.5), Inches(3.7), Inches(1.5),
         "La demande est là. La promesse de marque fonctionne.",
         size=12, italic=True, color=TEXTE)

add_rounded(s, Inches(4.85), y_big, Inches(4.0), Inches(2.8), BLANC,
            line_rgb=GRIS_LIGHT, radius=0.05)
add_rect(s, Inches(4.85), y_big, Inches(4.0), Pt(4), ROUGE)
add_text(s, Inches(5.15), Inches(3.6), Inches(3.7), Inches(0.4),
         "ABANDONS",
         size=10, bold=True, color=ROUGE)
add_text(s, Inches(5.15), Inches(4.0), Inches(3.7), Inches(1.0),
         "17,7 %",
         size=56, bold=True, color=NAVY, font="Georgia")
add_text(s, Inches(5.15), Inches(5.0), Inches(3.7), Inches(0.4),
         "1 142 apprenants perdus / an",
         size=12, color=GRIS)
add_text(s, Inches(5.15), Inches(5.5), Inches(3.7), Inches(1.5),
         "Stable depuis 12 mois. Le succès ne corrige pas le problème.",
         size=12, italic=True, color=TEXTE)

add_rounded(s, Inches(9.0), y_big, Inches(4.0), Inches(2.8), BLANC,
            line_rgb=GRIS_LIGHT, radius=0.05)
add_rect(s, Inches(9.0), y_big, Inches(4.0), Pt(4), ORANGE)
add_text(s, Inches(9.3), Inches(3.6), Inches(3.7), Inches(0.4),
         "ALERTES ML",
         size=10, bold=True, color=ORANGE)
add_text(s, Inches(9.3), Inches(4.0), Inches(3.7), Inches(1.0),
         "1 613",
         size=56, bold=True, color=NAVY, font="Georgia")
add_text(s, Inches(9.3), Inches(5.0), Inches(3.7), Inches(0.4),
         "77 % des apprenants actifs",
         size=12, color=GRIS)
add_text(s, Inches(9.3), Inches(5.5), Inches(3.7), Inches(1.5),
         "Le modèle voit venir une vague que les KPIs ne montrent pas encore.",
         size=12, italic=True, color=TEXTE)

add_footer(s, 2)


# ============================================================
# SLIDE 3 — 3 PARADOXES DATA-DRIVEN
# ============================================================

s = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s, 0, 0, SW, SH, FOND_PAGE)

add_text(s, Inches(0.7), Inches(0.5), Inches(3), Inches(0.3),
         "02  ·  CE QUE LA DATA RÉVÈLE",
         size=10, bold=True, color=VIOLET)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Trois paradoxes que les KPIs cachent",
         size=32, bold=True, color=NAVY, font="Georgia")

# 3 cartes de paradoxes
def paradox_card(slide, x, y, w, h, num, titre, fait, lecture, accent):
    add_rounded(slide, x, y, w, h, BLANC, line_rgb=GRIS_LIGHT, radius=0.04)
    add_rect(slide, x, y, w, Pt(4), accent)
    # Numéro grand
    add_text(slide, x + Inches(0.25), y + Inches(0.15), Inches(0.6), Inches(0.6),
             num, size=44, bold=True, color=accent, font="Georgia")
    # Titre
    add_text(slide, x + Inches(0.95), y + Inches(0.25), w - Inches(1.1), Inches(0.5),
             titre, size=14, bold=True, color=NAVY)
    # Fait chiffré
    add_text(slide, x + Inches(0.25), y + Inches(1.3), w - Inches(0.5), Inches(0.4),
             "LE FAIT", size=9, bold=True, color=accent)
    add_text(slide, x + Inches(0.25), y + Inches(1.65), w - Inches(0.5), Inches(0.9),
             fait, size=12, color=TEXTE)
    # Lecture
    add_text(slide, x + Inches(0.25), y + Inches(2.7), w - Inches(0.5), Inches(0.4),
             "LECTURE BUSINESS", size=9, bold=True, color=accent)
    add_text(slide, x + Inches(0.25), y + Inches(3.05), w - Inches(0.5), Inches(1.0),
             lecture, size=12, italic=True, color=NAVY, anchor=MSO_ANCHOR.TOP)

paradox_card(s, Inches(0.7), Inches(2.0), Inches(4.0), Inches(4.5),
    "01",
    "Satisfaits, mais ils partent",
    "CSAT 4,23/5 (élevé) ET 17,7 % d'abandon. Les apprenants apprécient le contenu mais ne le terminent pas.",
    "Ce n'est pas un problème de qualité — c'est un problème de friction. Charge, cadence, prérequis manquants.",
    VIOLET)

paradox_card(s, Inches(4.85), Inches(2.0), Inches(4.0), Inches(4.5),
    "02",
    "Concentration en haut, fragilité en bas",
    "Les Top 5 parcours portent 60 % du revenu. Mais 4 instructeurs sur 12 sont en zone « À coacher ».",
    "Si un instructeur clé part, jusqu'à 12 % du revenu est à risque. Le talent pédagogique est notre actif n°1.",
    ROUGE)

paradox_card(s, Inches(9.0), Inches(2.0), Inches(4.0), Inches(4.5),
    "03",
    "Les pires parcours sont les plus avancés",
    "Pire abandon par domaine : Machine Learning (16,9 %), Fullstack React (17,3 %), Growth Hacking (18,6 %), Mgmt Projet (18,4 %).",
    "Les apprenants surévaluent leur niveau au moment de l'inscription. Aucun système de prérequis ne filtre.",
    ORANGE)

add_footer(s, 3)


# ============================================================
# SLIDE 4 — LE MODÈLE ML
# ============================================================

s = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s, 0, 0, SW, SH, FOND_PAGE)

add_text(s, Inches(0.7), Inches(0.5), Inches(3), Inches(0.3),
         "03  ·  LE MODÈLE DE RISQUE",
         size=10, bold=True, color=VIOLET)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Sur 100 apprenants en cours, 77 sont à risque",
         size=28, bold=True, color=NAVY, font="Georgia")

add_text(s, Inches(0.7), Inches(1.85), Inches(12), Inches(0.4),
         "Le modèle ML segmente le portefeuille en 3 zones d'action et 1 zone neutre. Recall mesuré : 79,6 %.",
         size=13, italic=True, color=GRIS)

# 3 zones d'alerte
zones = [
    (ROUGE, ROUGE_LIGHT, "ROUGE", "Score > 0,80",   "~150",  "Intervention humaine",
     "Appel téléphonique du tuteur sous 48h.\nProbabilité d'abandon > 80 %."),
    (ORANGE, ORANGE_LIGHT, "ORANGE", "0,60 – 0,80", "~150",  "Relance ciblée",
     "Email personnalisé avec recommandation de parcours adapté.\nProbabilité 60–80 %."),
    (VIOLET, VIOLET_LIGHT, "VIOLET", "0,30 – 0,60", "~1 300", "Rappel doux",
     "Push notification de réengagement.\nProbabilité 30–60 %."),
]
zone_y = Inches(2.7)
for i, (col_main, col_light, label, score, n, action, detail) in enumerate(zones):
    x = Inches(0.7 + i * 4.15)
    w = Inches(4.0)
    add_rounded(s, x, zone_y, w, Inches(3.4), col_light, radius=0.04)
    add_rect(s, x, zone_y, w, Pt(4), col_main)
    # Header
    add_text(s, x + Inches(0.25), zone_y + Inches(0.2), Inches(2.0), Inches(0.4),
             label, size=14, bold=True, color=col_main)
    add_text(s, x + Inches(2.0), zone_y + Inches(0.2), Inches(1.8), Inches(0.4),
             score, size=10, color=GRIS, align=PP_ALIGN.RIGHT)
    # Big number
    add_text(s, x + Inches(0.25), zone_y + Inches(0.8), w - Inches(0.5), Inches(1.0),
             n, size=44, bold=True, color=NAVY, font="Georgia")
    add_text(s, x + Inches(0.25), zone_y + Inches(1.85), w - Inches(0.5), Inches(0.4),
             "apprenants concernés",
             size=11, color=GRIS, italic=True)
    # Action
    add_text(s, x + Inches(0.25), zone_y + Inches(2.35), w - Inches(0.5), Inches(0.4),
             action, size=13, bold=True, color=col_main)
    add_text(s, x + Inches(0.25), zone_y + Inches(2.75), w - Inches(0.5), Inches(0.6),
             detail, size=10, color=TEXTE)

# Note recall
add_rounded(s, Inches(0.7), Inches(6.3), Inches(12.3), Inches(0.55),
            VIOLET_LIGHT, radius=0.1)
add_text(s, Inches(1.0), Inches(6.4), Inches(12.0), Inches(0.4),
         "⚠ Recall 79,6 % : sur 100 apprenants qui décrocheront vraiment, le modèle en flag 80. La zone aveugle (20 %) doit être surveillée par un suivi qualitatif complémentaire.",
         size=11, italic=True, color=NAVY)

add_footer(s, 4)


# ============================================================
# SLIDE 5 — INSIGHTS NON ÉVIDENTS
# ============================================================

s = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s, 0, 0, SW, SH, FOND_PAGE)

add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3),
         "04  ·  CE QUE LA DATA NE CRIE PAS",
         size=10, bold=True, color=VIOLET)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Trois actifs business qu'on n'a pas vus venir",
         size=28, bold=True, color=NAVY, font="Georgia")

# 3 insights cachés
def hidden_insight(slide, x, y, w, h, num, titre, kpi, label, lecture, color):
    add_rounded(slide, x, y, w, h, BLANC, line_rgb=GRIS_LIGHT, radius=0.04)
    # Number badge
    badge = add_rounded(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.8), Inches(0.8),
                        color, radius=0.5)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.8), Inches(0.8),
             num, size=20, bold=True, color=BLANC, font="Georgia",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Titre
    add_text(slide, x + Inches(1.3), y + Inches(0.3), w - Inches(1.5), Inches(0.4),
             titre, size=14, bold=True, color=NAVY)
    add_text(slide, x + Inches(1.3), y + Inches(0.7), w - Inches(1.5), Inches(0.4),
             kpi + " · " + label, size=10, color=GRIS, italic=True)
    # Lecture
    add_text(slide, x + Inches(0.3), y + Inches(1.5), w - Inches(0.5), Inches(2.0),
             lecture, size=12, color=TEXTE, anchor=MSO_ANCHOR.TOP)

hidden_insight(s, Inches(0.7), Inches(2.0), Inches(4.0), Inches(4.5),
    "A",
    "Mobile Money = passeport régional",
    "50 %",
    "des paiements en Mobile Money",
    "Le marché ouest-africain est natif Mobile Money. Wave (Sénégal), Orange Money, MTN MoMo (Cameroun) : l'infra est déjà compatible.\n\nLe paiement n'est plus un frein à l'expansion — c'est notre tête de pont.",
    VERT)

hidden_insight(s, Inches(4.85), Inches(2.0), Inches(4.0), Inches(4.5),
    "B",
    "LTV apprenant = 173 000 FCFA",
    "264 €",
    "revenu généré par apprenant actif",
    "778 M FCFA / 4 500 apprenants. Et seulement 16 % d'acquisition payante.\n\nLa LTV est saine, le coût d'acquisition est faible : il y a de la place pour scaler le marketing payant sans détruire l'unit economics.",
    VIOLET)

hidden_insight(s, Inches(9.0), Inches(2.0), Inches(4.0), Inches(4.5),
    "C",
    "Le SEO est un canal sous-exploité",
    "35 %",
    "des apprenants viennent de l'organique",
    "Le bouche-à-oreille et le SEO ramènent 1 apprenant sur 3, gratuitement. C'est notre canal le plus rentable.\n\nDoubler l'investissement contenu (blog, témoignages, SEO technique) peut amplifier ce levier sans gonfler le CAC.",
    ORANGE)

add_footer(s, 5)


# ============================================================
# SLIDE 6 — RECOMMANDATION 1
# ============================================================

def reco_slide(num, slide_num, color, color_light, titre, hook, action_lines,
               impact_kpi, impact_label, delai, cout):
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_rect(s, 0, 0, SW, SH, FOND_PAGE)

    add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3),
             f"05  ·  RECOMMANDATION {num}",
             size=10, bold=True, color=color)
    add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
             titre, size=32, bold=True, color=NAVY, font="Georgia")

    # Hook
    add_rounded(s, Inches(0.7), Inches(2.0), Inches(12.3), Inches(0.7),
                color_light, radius=0.1)
    add_text(s, Inches(1.0), Inches(2.1), Inches(11.8), Inches(0.5),
             hook, size=14, italic=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)

    # Section actions (gauche) + impact (droite)
    add_text(s, Inches(0.7), Inches(3.1), Inches(0.5), Inches(0.4),
             "ACTION", size=10, bold=True, color=color)
    actions_y = Inches(3.55)
    for i, (head, body) in enumerate(action_lines):
        y = actions_y + Inches(i * 0.85)
        # bullet point
        add_rect(s, Inches(0.75), y + Inches(0.12), Inches(0.08), Inches(0.08), color)
        add_text(s, Inches(0.95), y, Inches(7.5), Inches(0.35),
                 head, size=13, bold=True, color=NAVY)
        add_text(s, Inches(0.95), y + Inches(0.35), Inches(7.5), Inches(0.5),
                 body, size=11, color=GRIS, italic=False)

    # Impact card (droite)
    add_rounded(s, Inches(8.7), Inches(3.1), Inches(4.3), Inches(3.6),
                NAVY_DARK, radius=0.04)
    add_text(s, Inches(8.95), Inches(3.3), Inches(3.8), Inches(0.4),
             "IMPACT ATTENDU", size=10, bold=True,
             color=RGBColor(0x9A, 0xB6, 0xE0))
    add_text(s, Inches(8.95), Inches(3.7), Inches(3.8), Inches(1.4),
             impact_kpi, size=44, bold=True, color=BLANC, font="Georgia")
    add_text(s, Inches(8.95), Inches(5.05), Inches(3.8), Inches(0.5),
             impact_label, size=12, italic=True,
             color=RGBColor(0xCB, 0xD5, 0xE0))
    # Sep
    add_rect(s, Inches(8.95), Inches(5.7), Inches(0.5), Pt(2),
             RGBColor(0x4A, 0x5F, 0x88))
    add_text(s, Inches(8.95), Inches(5.85), Inches(3.8), Inches(0.4),
             "Délai", size=9, bold=True,
             color=RGBColor(0x9A, 0xB6, 0xE0))
    add_text(s, Inches(8.95), Inches(6.1), Inches(3.8), Inches(0.4),
             delai, size=12, color=BLANC)
    add_text(s, Inches(11), Inches(5.85), Inches(2), Inches(0.4),
             "Investissement", size=9, bold=True,
             color=RGBColor(0x9A, 0xB6, 0xE0))
    add_text(s, Inches(11), Inches(6.1), Inches(2), Inches(0.4),
             cout, size=12, color=BLANC)

    add_footer(s, slide_num)
    return s


reco_slide(
    "01", 6, VIOLET, VIOLET_LIGHT,
    "Réduire la friction d'entrée sur les parcours avancés",
    "Les apprenants quittent ML, Fullstack et Growth car ils sous-estiment la marche d'entrée. Pas par manque de qualité.",
    [
        ("Test de prérequis obligatoire avant inscription",
         "10 questions techniques + 5 questions de motivation. Score < seuil → orientation vers le parcours fondations."),
        ("Onboarding séquencé sur 5 jours",
         "Mini-modules quotidiens pour acclimater au rythme. CTA quotidien de 15 minutes."),
        ("Buddy system entre apprenants",
         "Pairer chaque nouvel inscrit avec un alumni du même parcours. 1 contact / semaine pendant 4 semaines."),
        ("Refonte du syllabus des 4 pires parcours",
         "Découper les chapitres lourds, ajouter des micro-livrables, retirer les contenus à faible CSAT identifiés."),
    ],
    "+5 pp",
    "complétion sur les 4 parcours ciblés (33,3 % → 38 %)",
    "6 mois",
    "12 M FCFA",
)


# ============================================================
# SLIDE 7 — RECOMMANDATION 2
# ============================================================

reco_slide(
    "02", 7, ROUGE, ROUGE_LIGHT,
    "Activer une cellule de rétention 3 niveaux",
    "Le modèle ML identifie déjà les apprenants à risque. Il manque le bras opérationnel qui agit sur le signal.",
    [
        ("Niveau ROUGE — Cellule humaine de tutorat (~150 apprenants)",
         "1 tuteur dédié, appel sous 48h, plan de remédiation personnalisé sur 2 semaines."),
        ("Niveau ORANGE — Email perso + visio collective (~150 apprenants)",
         "Recommandation de parcours adapté ou bascule vers un parcours fondations. Visio mensuelle de cohorte."),
        ("Niveau VIOLET — Push automatisé contextuel (~1 300 apprenants)",
         "Notifications déclenchées par le score ML : J+3 inactivité, J+7, J+14, avec contenu pertinent."),
        ("Mesure obsessionnelle de la conversion",
         "Re-engagement à 7 et 30 jours. Si une zone ne convertit pas, on coupe et on réalloue le budget."),
    ],
    "−4 pp",
    "abandon (17,7 % → 13,7 %)",
    "3 mois",
    "8 M FCFA",
)


# ============================================================
# SLIDE 8 — RECOMMANDATION 3
# ============================================================

reco_slide(
    "03", 8, VERT, VERT_LIGHT,
    "Passer de 1,44 à 2,5 inscriptions par apprenant",
    "Le cross-sell est quasi nul. Doubler la LTV ne demande pas de nouveaux apprenants — il suffit de mieux servir les existants.",
    [
        ("Bundle de parcours par carrière (4 packs)",
         "Data Analyst (Python + SQL + Power BI), Dev Web (HTML/CSS/JS + React + Node), Growth (SEO + Content + Ads), Manager (Leadership + Mgmt Projet + Soft Skills)."),
        ("Réduction de 25 % sur le 2ᵉ parcours",
         "Activée automatiquement à la complétion d'un premier parcours. Mobile Money en facilitateur (paiement en 3 fois)."),
        ("Recommandations IA personnalisées",
         "À la complétion, suggérer 2 parcours complémentaires basés sur le profil + similarités d'apprenants ayant réussi."),
        ("Programme alumni avec parcours masterclasses",
         "Accès payant à des masterclasses trimestrielles animées par les Top Performers. Récurrence revenu."),
    ],
    "+220M FCFA",
    "+28 % de revenu (778M → ~1 Md FCFA)",
    "6 mois",
    "5 M FCFA",
)


# ============================================================
# SLIDE 9 — ROADMAP
# ============================================================

s = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s, 0, 0, SW, SH, FOND_PAGE)

add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3),
         "06  ·  EXÉCUTION",
         size=10, bold=True, color=VIOLET)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Roadmap 6 mois — focus sur les jalons mesurables",
         size=28, bold=True, color=NAVY, font="Georgia")

# Timeline horizontale
tl_y = Inches(2.6)
tl_x = Inches(0.7)
tl_w = Inches(12.0)
# Ligne
add_rect(s, tl_x, tl_y + Inches(0.15), tl_w, Pt(2), GRIS_LIGHT)
# 3 jalons
jalons = [
    (Inches(0.7),  "M+1", "Lancement",
     "Cellule ROUGE active.\nTest de prérequis en ligne.\nPremiers bundles vendus.",
     VERT),
    (Inches(5.7),  "M+3", "Premier bilan",
     "ORANGE et VIOLET activés.\nRevue pédagogique des 4 instructeurs.\nMesure abandon trimestriel.",
     ORANGE),
    (Inches(10.7), "M+6", "Évaluation finale",
     "Cible complétion +5 pp.\nCible abandon −4 pp.\nCible revenu +220M FCFA.",
     VIOLET),
]
for x, m, titre, body, color in jalons:
    # Cercle de jalon
    cir = s.shapes.add_shape(MSO_SHAPE.OVAL, x, tl_y - Inches(0.05), Inches(0.5), Inches(0.5))
    cir.fill.solid(); cir.fill.fore_color.rgb = color
    cir.line.fill.background(); cir.shadow.inherit = False
    # Mois en blanc dans le cercle
    add_text(s, x, tl_y - Inches(0.05), Inches(0.5), Inches(0.5),
             m, size=10, bold=True, color=BLANC, font="Calibri",
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Carte sous le jalon
    add_rounded(s, x - Inches(1.2), Inches(3.4), Inches(2.9), Inches(2.7),
                BLANC, line_rgb=GRIS_LIGHT, radius=0.04)
    add_rect(s, x - Inches(1.2), Inches(3.4), Inches(2.9), Pt(4), color)
    add_text(s, x - Inches(1.0), Inches(3.6), Inches(2.6), Inches(0.4),
             titre, size=15, bold=True, color=NAVY, font="Georgia")
    add_text(s, x - Inches(1.0), Inches(4.1), Inches(2.6), Inches(2.0),
             body, size=11, color=TEXTE)

# Total impact en bas
add_rounded(s, Inches(0.7), Inches(6.4), Inches(12.3), Inches(0.7),
            NAVY, radius=0.05)
add_text(s, Inches(1.0), Inches(6.45), Inches(12.0), Inches(0.6),
         "Impact agrégé à 6 mois : +5 pp complétion · −4 pp abandon · +220 M FCFA de revenu pour 25 M FCFA d'investissement total. ROI ≈ x9.",
         size=13, bold=True, color=BLANC,
         anchor=MSO_ANCHOR.MIDDLE)

add_footer(s, 9)


# ============================================================
# SLIDE 10 — LIMITES & PROCHAINES DÉCISIONS
# ============================================================

s = prs.slides.add_slide(BLANK_LAYOUT)
add_rect(s, 0, 0, SW, SH, FOND_PAGE)

add_text(s, Inches(0.7), Inches(0.5), Inches(4), Inches(0.3),
         "07  ·  CONDITIONS DE RÉUSSITE",
         size=10, bold=True, color=VIOLET)
add_text(s, Inches(0.7), Inches(0.9), Inches(12), Inches(1.0),
         "Ce qu'il faut décider, ce qu'on ne sait pas encore",
         size=28, bold=True, color=NAVY, font="Georgia")

# Colonne gauche : limites du modèle
add_text(s, Inches(0.7), Inches(2.1), Inches(6), Inches(0.4),
         "LIMITES DU MODÈLE & DES DONNÉES", size=11, bold=True, color=ROUGE)

limites = [
    ("Recall 79,6 %",
     "20 % des futurs abandons ne sont pas détectés. À compenser par une enquête qualitative à 30 jours."),
    ("Saisonnalité partiellement modélisée",
     "Le pic de juin est expliqué à 60 % seulement. Le reste tient à des facteurs externes non capturés."),
    ("Pas de feedback post-abandon",
     "On ne sait pas pourquoi un apprenant arrête. Mettre en place un mini-questionnaire à J+30 après abandon."),
    ("Réentraînement trimestriel nécessaire",
     "Les comportements évoluent. Sans réentraînement, le recall chutera de 5 pp tous les 6 mois."),
]
for i, (head, body) in enumerate(limites):
    y = Inches(2.55 + i * 1.05)
    add_rect(s, Inches(0.7), y + Inches(0.1), Pt(3), Inches(0.85), ROUGE)
    add_text(s, Inches(0.85), y, Inches(5.8), Inches(0.4),
             head, size=12, bold=True, color=NAVY)
    add_text(s, Inches(0.85), y + Inches(0.3), Inches(5.8), Inches(0.7),
             body, size=10, color=GRIS)

# Colonne droite : décisions à prendre
add_text(s, Inches(7.0), Inches(2.1), Inches(6), Inches(0.4),
         "DÉCISIONS ATTENDUES DU COMITÉ", size=11, bold=True, color=VIOLET)

decisions = [
    ("Validation du budget 25 M FCFA",
     "Ouverture des 3 chantiers en parallèle ou séquencé. Préférence : parallèle (économies d'échelle sur le coaching)."),
    ("Désignation d'un sponsor",
     "Membre du COMEX porteur des 3 recommandations. Reporting mensuel obligatoire."),
    ("Choix du KPI nord",
     "Complétion (orienté pédago) ou Revenu/apprenant (orienté business) ? Une seule étoile polaire pour l'équipe."),
    ("Communication interne",
     "Annonce aux instructeurs « À coacher » : posture d'accompagnement, pas de sanction. Tonalité critique."),
]
for i, (head, body) in enumerate(decisions):
    y = Inches(2.55 + i * 1.05)
    add_rect(s, Inches(7.0), y + Inches(0.1), Pt(3), Inches(0.85), VIOLET)
    add_text(s, Inches(7.15), y, Inches(5.8), Inches(0.4),
             head, size=12, bold=True, color=NAVY)
    add_text(s, Inches(7.15), y + Inches(0.3), Inches(5.8), Inches(0.7),
             body, size=10, color=GRIS)

add_footer(s, 10)


# ============================================================
# Sauvegarde
# ============================================================

prs.save(str(OUT))
print(f"\n[OK] Synthese generee : {OUT}")
print(f"     Taille : {OUT.stat().st_size / 1024:.1f} ko")
print(f"     Slides : 10")
